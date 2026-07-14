from __future__ import annotations

import json
import hashlib
import os
import re
import shutil
import zipfile
import uuid
from datetime import datetime
from pathlib import Path
from time import perf_counter, time
from typing import Any

from flask import Flask, jsonify, render_template, request, send_file
from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
JOBS_DIR = BASE_DIR / "_fluoro_jobs"
EXPORT_DIR = BASE_DIR / "exports"
SUPPORTED_EXTENSIONS = {".tif", ".tiff", ".png", ".jpg", ".jpeg", ".bmp"}
MAX_PREVIEW_SIZE = (900, 650)
RASTER_DPI = 240
TEMP_TTL_SECONDS = 24 * 60 * 60
JOB_ID_PATTERN = re.compile(r"^[0-9a-f]{32}$")
PREVIEW_NAME_PATTERN = re.compile(r"^\d{5}_[0-9a-f]{8}\.png$")
EXPORT_PROFILES = {
    "nature-single": (8.9, None),
    "nature-double": (18.3, None),
    "slides": (33.866, 19.05),
}

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 8 * 1024 * 1024 * 1024


def now_stamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def unique_export_stamp() -> str:
    return f"{now_stamp()}_{uuid.uuid4().hex[:12]}"


def validated_job_id(job_id: Any) -> str:
    value = str(job_id or "")
    if not JOB_ID_PATTERN.fullmatch(value):
        raise FileNotFoundError("任务不存在或已过期，请重新上传")
    return value


def job_export_dir(manifest: dict[str, Any]) -> Path:
    output_dir = EXPORT_DIR / validated_job_id(manifest.get("job_id"))
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir


def cleanup_expired_temp_data() -> None:
    cutoff = time() - TEMP_TTL_SECONDS
    for root in (JOBS_DIR, EXPORT_DIR):
        if not root.exists():
            continue
        for child in root.iterdir():
            try:
                if child.stat().st_mtime >= cutoff:
                    continue
                if child.is_dir():
                    shutil.rmtree(child, ignore_errors=True)
                else:
                    child.unlink(missing_ok=True)
            except OSError:
                continue


def safe_path(raw_path: str) -> Path:
    normalized = raw_path.replace("\\", "/")
    parts = []
    for part in normalized.split("/"):
        part = part.strip()
        if not part or part in {".", ".."}:
            continue
        part = re.sub(r'[<>:"|?*\x00-\x1f]', "_", part)
        parts.append(part[:140] or "file")
    return Path(*parts) if parts else Path(f"image_{uuid.uuid4().hex[:8]}")


def is_supported_image(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as source:
        for chunk in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def stable_image_id(image: dict[str, Any], index: int = 0) -> str:
    existing = str(image.get("image_id") or "").strip()
    if existing:
        return existing
    identity = f"{image.get('relative_path', '')}\0{image.get('filename', '')}\0{index}"
    return hashlib.sha256(identity.encode("utf-8")).hexdigest()[:24]


def ordered_manifest_images(manifest: dict[str, Any], settings: dict[str, Any] | None = None) -> list[dict[str, Any]]:
    images = []
    for index, raw in enumerate(manifest.get("images", [])):
        image = dict(raw)
        image["image_id"] = stable_image_id(image, index)
        images.append(image)
    requested_order = (settings or {}).get("image_order")
    if not isinstance(requested_order, list):
        return images
    rank = {str(image_id): index for index, image_id in enumerate(requested_order)}
    original_rank = {image["image_id"]: index for index, image in enumerate(images)}
    return sorted(images, key=lambda image: (rank.get(image["image_id"], len(rank) + original_rank[image["image_id"]])))


def resolve_export_profile(settings: dict[str, Any]) -> dict[str, Any]:
    name = str(settings.get("export_profile") or "slides").strip().lower()
    if name == "custom":
        try:
            width_cm = float(settings.get("export_width_cm"))
            height_cm = float(settings.get("export_height_cm"))
        except (TypeError, ValueError):
            raise ValueError("自定义导出尺寸需要填写宽度和高度（cm）") from None
    elif name in EXPORT_PROFILES:
        width_cm, height_cm = EXPORT_PROFILES[name]
        if height_cm is None:
            height_cm = width_cm * 7.5 / 13.333
    else:
        raise ValueError("导出规格仅支持 nature-single、nature-double、custom 或 slides")

    if not 3 <= width_cm <= 40 or not 2 <= height_cm <= 40:
        raise ValueError("导出尺寸需在宽 3–40 cm、高 2–40 cm 范围内")
    try:
        dpi = int(settings.get("export_dpi") or 300)
    except (TypeError, ValueError):
        raise ValueError("导出分辨率必须为 300 或 600 DPI") from None
    if dpi not in {300, 600}:
        raise ValueError("导出分辨率必须为 300 或 600 DPI")

    width_in = width_cm / 2.54
    height_in = height_cm / 2.54
    logical_width = 13.333
    logical_height = height_in * logical_width / width_in
    return {
        "name": name,
        "width_cm": round(width_cm, 4),
        "height_cm": round(height_cm, 4),
        "width_in": width_in,
        "height_in": height_in,
        "dpi": dpi,
        "pixel_width": max(1, round(width_in * dpi)),
        "pixel_height": max(1, round(height_in * dpi)),
        "logical_width": logical_width,
        "logical_height": logical_height,
        "coordinate_scale": width_in / logical_width,
    }


def normalized_stem(filename: str) -> str:
    name = Path(filename).name
    for _ in range(4):
        before = name
        name = re.sub(r"(?i)\.(tif|tiff|png|jpe?g|bmp)$", "", name)
        name = re.sub(r"(?i)\.(tif|tiff|png|jpe?g|bmp)-\d+$", "", name)
        if name == before:
            break
    return name


def detect_channel(filename: str) -> dict[str, Any]:
    stem = normalized_stem(filename)
    lower = stem.lower()

    if re.search(r"(?i)(^|[_\-\s])(merged?|merge)(?=$|[_\-\s])", stem):
        return {"key": "Merged", "label": "Merge", "order": 10000}

    channel_match = re.search(r"(?i)(^|[_\-\s])ch(?:annel)?[_\-\s]*0*(\d+)(?=$|[_\-\s])", stem)
    if channel_match:
        channel_number = int(channel_match.group(2))
        return {
            "key": f"ch{channel_number:02d}",
            "label": f"ch{channel_number:02d}",
            "order": channel_number,
        }

    known_markers = [
        ("dapi", "DAPI", 0),
        ("hoechst", "Hoechst", 0),
        ("gfp", "GFP", 1),
        ("fitc", "FITC", 1),
        ("cy3", "Cy3", 2),
        ("tritc", "TRITC", 2),
        ("rfp", "RFP", 2),
        ("cy5", "Cy5", 3),
    ]
    for token, label, order in known_markers:
        if re.search(rf"(?i)(^|[_\-\s]){re.escape(token)}(?=$|[_\-\s])", lower):
            return {"key": label, "label": label, "order": order}

    return {"key": "Image", "label": "Image", "order": 5000}


def sample_key_from_path(relative_path: str) -> str:
    path = Path(relative_path.replace("\\", "/"))
    stem = normalized_stem(path.name)
    stem = re.sub(r"(?i)^merged?[_\-\s]*", "", stem)
    stem = re.sub(r"(?i)(^|[_\-\s])merged?(?=$|[_\-\s])", "_", stem)
    stem = re.sub(r"(?i)(^|[_\-\s])merge(?=$|[_\-\s])", "_", stem)
    stem = re.sub(r"(?i)(^|[_\-\s])ch(?:annel)?[_\-\s]*0*\d+(?=$|[_\-\s])", "_", stem)
    stem = re.sub(r"[_\-\s]+", "_", stem).strip("_- ")
    stem = stem or normalized_stem(path.name)
    parent = path.parent.as_posix()
    return f"{parent}/{stem}" if parent and parent != "." else stem


def channel_sort_key(channel: str) -> tuple[int, str]:
    if channel.lower() == "merged":
        return (10000, channel)
    match = re.match(r"(?i)^ch0*(\d+)$", channel)
    if match:
        return (int(match.group(1)), channel)
    return (5000, channel.lower())


def image_to_rgb(image: Image.Image) -> Image.Image:
    image = ImageOps.exif_transpose(image)
    if image.mode == "RGBA":
        background = Image.new("RGBA", image.size, (0, 0, 0, 255))
        return Image.alpha_composite(background, image).convert("RGB")
    if image.mode in {"RGB", "L"}:
        return image.convert("RGB")
    if image.mode.startswith("I;16"):
        # Pillow's direct I;16 -> L conversion clips every value above 255.
        # Use the full detector range so 16-bit microscopy intensities remain ordered.
        return image.convert("I").point(lambda value: value * 255 / 65535).convert("L").convert("RGB")
    if image.mode in {"I", "F"}:
        low, high = image.getextrema()
        if image.mode == "I" and low >= 0 and high <= 65535:
            low, high = 0, 65535
        elif image.mode == "F" and low >= 0 and high <= 1:
            low, high = 0.0, 1.0
        if high <= low:
            return Image.new("RGB", image.size, (0, 0, 0))
        scale = 255 / (high - low)
        return image.point(lambda value: (value - low) * scale).convert("L").convert("RGB")
    return image.convert("RGB")


def prepare_image(source_path: Path, preview_path: Path) -> dict[str, Any]:
    with Image.open(source_path) as source:
        source.seek(0)
        original_format = source.format or source_path.suffix.lstrip(".").upper()
        original_mode = source.mode
        frame_count = int(getattr(source, "n_frames", 1))
        dpi_value = source.info.get("dpi")
        tags = getattr(source, "tag_v2", None)
        image = ImageOps.exif_transpose(source)
        width, height = image.size
        if isinstance(dpi_value, (tuple, list)):
            resolution_dpi = [round(float(value), 4) for value in dpi_value[:2]]
        elif isinstance(dpi_value, (int, float)):
            resolution_dpi = [round(float(dpi_value), 4)] * 2
        else:
            resolution_dpi = None
        bits_per_sample = None
        if tags:
            raw_bits = tags.get(258)
            if isinstance(raw_bits, (tuple, list)):
                bits_per_sample = [int(value) for value in raw_bits]
            elif raw_bits is not None:
                bits_per_sample = [int(raw_bits)]
        metadata = {
            "format": original_format,
            "mime_type": Image.MIME.get(original_format),
            "mode": original_mode,
            "frame_count": frame_count,
            "resolution_dpi": resolution_dpi,
            "bits_per_sample": bits_per_sample,
            "display_mapping": (
                "linear_full_16bit_to_8bit"
                if original_mode.startswith("I;16")
                else "linear_numeric_range_to_8bit" if original_mode in {"I", "F"} else "native_8bit_or_rgb"
            ),
        }
        preview_path.parent.mkdir(parents=True, exist_ok=True)
        preview = image.copy()
        preview.thumbnail(MAX_PREVIEW_SIZE, Image.Resampling.LANCZOS)
        preview = image_to_rgb(preview)
        preview.save(preview_path, "PNG", compress_level=4)
    return {"width": width, "height": height, "metadata": metadata}


def make_job_manifest(job_id: str, items: list[dict[str, Any]]) -> dict[str, Any]:
    groups_by_key: dict[str, dict[str, Any]] = {}
    channels: dict[str, dict[str, Any]] = {}
    ordered_images: list[dict[str, Any]] = []

    for item in items:
        item = dict(item)
        item["image_id"] = stable_image_id(item, len(ordered_images))
        group_key = sample_key_from_path(item["relative_path"])
        channel_info = detect_channel(item["relative_path"])
        base_channel = channel_info["key"]
        group = groups_by_key.setdefault(
            group_key,
            {"key": group_key, "display": group_key, "images": []},
        )

        channel_key = base_channel
        used = {image["channel"] for image in group["images"]}
        duplicate_index = 2
        while channel_key in used:
            channel_key = f"{base_channel}_{duplicate_index}"
            duplicate_index += 1

        image_item = {
            **item,
            "channel": channel_key,
            "channel_label": channel_info["label"] if channel_key == base_channel else channel_key,
            "channel_order": channel_info["order"],
        }
        group["images"].append(image_item)
        ordered_images.append(image_item)
        channels.setdefault(
            channel_key,
            {
                "key": channel_key,
                "label": image_item["channel_label"],
                "order": channel_info["order"],
            },
        )

    ordered_groups = sorted(groups_by_key.values(), key=lambda group: group["key"].lower())
    for group in ordered_groups:
        group["images"].sort(key=lambda image: (image["channel_order"], image["channel"].lower()))

    ordered_channels = sorted(channels.values(), key=lambda item: (item["order"], item["key"].lower()))

    return {
        "job_id": job_id,
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "images": ordered_images,
        "groups": ordered_groups,
        "channels": ordered_channels,
        "image_count": len(items),
    }


def manifest_path(job_id: str) -> Path:
    return JOBS_DIR / validated_job_id(job_id) / "manifest.json"


def load_manifest(job_id: str) -> dict[str, Any]:
    path = manifest_path(job_id)
    if not path.exists():
        raise FileNotFoundError("任务不存在或已过期，请重新上传")
    return json.loads(path.read_text(encoding="utf-8"))


def public_manifest(manifest: dict[str, Any]) -> dict[str, Any]:
    image_ids = {
        str(image.get("relative_path") or image.get("filename") or index): stable_image_id(image, index)
        for index, image in enumerate(manifest.get("images", []))
    }
    result = {
        "job_id": manifest["job_id"],
        "created_at": manifest["created_at"],
        "image_count": manifest["image_count"],
        "images": [],
        "channels": manifest["channels"],
        "groups": [],
    }
    for index, image in enumerate(manifest.get("images", [])):
        image_id = stable_image_id(image, index)
        result["images"].append(
            {
                "image_id": image_id,
                "relative_path": image["relative_path"],
                "filename": image["filename"],
                "width": image["width"],
                "height": image["height"],
                "sha256": image.get("sha256"),
                "byte_size": image.get("byte_size"),
                "original_metadata": image.get("original_metadata") or {},
                "preview_url": f"/api/jobs/{manifest['job_id']}/preview/{image['preview_name']}",
            }
        )
    for group in manifest["groups"]:
        result_group = {"key": group["key"], "display": group["display"], "images": []}
        for image in group["images"]:
            image_identity = str(image.get("relative_path") or image.get("filename") or "")
            result_group["images"].append(
                {
                    "image_id": image.get("image_id") or image_ids.get(image_identity),
                    "relative_path": image["relative_path"],
                    "filename": image["filename"],
                    "channel": image["channel"],
                    "channel_label": image["channel_label"],
                    "width": image["width"],
                    "height": image["height"],
                    "sha256": image.get("sha256"),
                    "byte_size": image.get("byte_size"),
                    "original_metadata": image.get("original_metadata") or {},
                    "preview_url": f"/api/jobs/{manifest['job_id']}/preview/{image['preview_name']}",
                }
            )
        result["groups"].append(result_group)
    return result


def add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    rotation: int | None = None,
    scale: float = 1.0,
) -> None:
    shape = slide.shapes.add_textbox(Inches(left * scale), Inches(top * scale), Inches(width * scale), Inches(height * scale))
    if rotation is not None:
        shape.rotation = rotation
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(max(7, font_size * scale) if scale < 1 else font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def center_crop_image(source_path: Path, output_path: Path, target_aspect: float) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with Image.open(source_path) as source:
        rgb = image_to_rgb(source)
        width, height = rgb.size
        source_aspect = width / height
        if source_aspect > target_aspect:
            new_width = int(height * target_aspect)
            left = max(0, (width - new_width) // 2)
            box = (left, 0, left + new_width, height)
        else:
            new_height = int(width / target_aspect)
            top = max(0, (height - new_height) // 2)
            box = (0, top, width, top + new_height)
        rgb.crop(box).save(output_path, "PNG", compress_level=4)
    return output_path


def image_source_path(image: dict[str, Any]) -> Path:
    return Path(image.get("converted_path") or image["source_path"])


def calibration_for_image(settings: dict[str, Any], image: dict[str, Any] | None = None) -> dict[str, Any]:
    calibrations = settings.get("calibrations") or {}
    if not isinstance(calibrations, dict):
        raise ValueError("比例尺校准参数格式无效")
    if image:
        by_image = calibrations.get("by_image") or {}
        by_channel = calibrations.get("by_channel") or {}
        if not isinstance(by_image, dict) or not isinstance(by_channel, dict):
            raise ValueError("按图片或通道的比例尺校准格式无效")
        image_id = str(image.get("image_id") or "")
        channel = str(image.get("channel") or "")
        if image_id and isinstance(by_image.get(image_id), dict):
            return by_image[image_id]
        if channel and isinstance(by_channel.get(channel), dict):
            return by_channel[channel]
        if str(settings.get("figure_type") or "if").lower() == "ihc" and channel in {"low", "high"}:
            raise ValueError(f"IHC {channel} 倍率缺少独立的像素/微米校准")
    default = calibrations.get("default")
    if isinstance(default, dict):
        return default
    legacy = settings.get("calibration") or {}
    if not isinstance(legacy, dict):
        raise ValueError("比例尺校准参数格式无效")
    return legacy


def resolve_scale_bar(settings: dict[str, Any], image: dict[str, Any] | None = None) -> dict[str, Any] | None:
    raw = settings.get("scale_bar")
    if not isinstance(raw, dict) or not raw.get("enabled"):
        return None
    calibration = calibration_for_image(settings, image)
    try:
        length_um = float(raw.get("length_um"))
        pixels_per_micron = float(calibration.get("pixels_per_micron") or 0)
        if pixels_per_micron <= 0:
            known_px = float(calibration.get("known_distance_px") or 0)
            known_um = float(calibration.get("known_distance_um") or 0)
            pixels_per_micron = known_px / known_um if known_px > 0 and known_um > 0 else 0
        thickness_px = float(raw.get("thickness_px") or 4)
    except (TypeError, ValueError, ZeroDivisionError):
        raise ValueError("比例尺长度和校准值必须为有效数字") from None
    if length_um <= 0 or pixels_per_micron <= 0:
        raise ValueError("启用比例尺后需填写正数长度和像素/微米校准值")
    if not 1 <= thickness_px <= 30:
        raise ValueError("比例尺线宽需在 1–30 px 之间")
    position = str(raw.get("position") or "bottom-right").lower()
    if position not in {"bottom-left", "bottom-right", "top-left", "top-right"}:
        raise ValueError("比例尺位置无效")
    color = str(raw.get("color") or "white").lower()
    if color not in {"white", "black"}:
        raise ValueError("比例尺颜色仅支持 white 或 black")
    return {
        "length_um": length_um,
        "pixels_per_micron": pixels_per_micron,
        "source_length_px": length_um * pixels_per_micron,
        "position": position,
        "color": color,
        "thickness_px": thickness_px,
        "label": f"{length_um:g} μm",
    }


def image_display_geometry(
    source_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    fit_mode: str,
) -> tuple[float, float, float, float, float]:
    with Image.open(source_path) as image:
        source_width, source_height = image.size
    source_aspect = source_width / source_height
    cell_aspect = width / height
    if fit_mode == "contain":
        if source_aspect >= cell_aspect:
            draw_width = width
            draw_height = width / source_aspect
        else:
            draw_height = height
            draw_width = height * source_aspect
        return (
            left + (width - draw_width) / 2,
            top + (height - draw_height) / 2,
            draw_width,
            draw_height,
            float(source_width),
        )
    visible_source_width = min(float(source_width), float(source_height) * cell_aspect)
    return left, top, width, height, visible_source_width


def add_scale_bar_pptx(
    slide,
    source_path: Path,
    box: tuple[float, float, float, float],
    fit_mode: str,
    settings: dict[str, Any],
    coordinate_scale: float,
    image: dict[str, Any] | None = None,
) -> None:
    scale_bar = resolve_scale_bar(settings, image)
    if not scale_bar:
        return
    left, top, width, height, visible_source_width = image_display_geometry(source_path, *box, fit_mode)
    bar_width = width * scale_bar["source_length_px"] / visible_source_width
    if bar_width > width * 0.9:
        raise ValueError(f"比例尺 {scale_bar['label']} 超出图片可见宽度")
    margin = min(width, height) * 0.06
    on_right = scale_bar["position"].endswith("right")
    on_bottom = scale_bar["position"].startswith("bottom")
    start_x = left + width - margin - bar_width if on_right else left + margin
    end_x = start_x + bar_width
    y = top + height - margin if on_bottom else top + margin
    color = RGBColor(255, 255, 255) if scale_bar["color"] == "white" else RGBColor(0, 0, 0)
    add_ppt_line(
        slide,
        (start_x, y),
        (end_x, y),
        color,
        coordinate_scale,
        max(0.75, scale_bar["thickness_px"] * 72 / 300),
    )
    label_height = 0.20 / max(coordinate_scale, 0.1)
    label_top = y - label_height if on_bottom else y
    add_textbox(
        slide,
        scale_bar["label"],
        start_x - margin,
        label_top,
        bar_width + margin * 2,
        label_height,
        8,
        color,
        bold=False,
        align=PP_ALIGN.CENTER,
        scale=coordinate_scale,
    )


def draw_scale_bar_raster(
    canvas: Image.Image,
    source_path: Path,
    box: tuple[int, int, int, int],
    fit_mode: str,
    settings: dict[str, Any],
    dpi: int,
    image: dict[str, Any] | None = None,
) -> None:
    scale_bar = resolve_scale_bar(settings, image)
    if not scale_bar:
        return
    left, top, width, height, visible_source_width = image_display_geometry(
        source_path,
        float(box[0]),
        float(box[1]),
        float(box[2] - box[0]),
        float(box[3] - box[1]),
        fit_mode,
    )
    bar_width = width * scale_bar["source_length_px"] / visible_source_width
    if bar_width > width * 0.9:
        raise ValueError(f"比例尺 {scale_bar['label']} 超出图片可见宽度")
    margin = max(4, round(min(width, height) * 0.06))
    on_right = scale_bar["position"].endswith("right")
    on_bottom = scale_bar["position"].startswith("bottom")
    start_x = round(left + width - margin - bar_width if on_right else left + margin)
    end_x = round(start_x + bar_width)
    y = round(top + height - margin if on_bottom else top + margin)
    color = (255, 255, 255) if scale_bar["color"] == "white" else (0, 0, 0)
    thickness = max(1, round(scale_bar["thickness_px"] * dpi / 300))
    draw = ImageDraw.Draw(canvas)
    draw.line((start_x, y, end_x, y), fill=color, width=thickness)
    typeface = font(max(8, round(8 * dpi / 72)))
    text_box = draw.textbbox((0, 0), scale_bar["label"], font=typeface)
    text_width = text_box[2] - text_box[0]
    text_height = text_box[3] - text_box[1]
    text_x = (start_x + end_x - text_width) / 2
    text_y = y - text_height - thickness * 2 - text_box[1] if on_bottom else y + thickness * 2 - text_box[1]
    draw.text((text_x, text_y), scale_bar["label"], fill=color, font=typeface)


def add_image_cover(
    slide,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    crop_dir: Path,
    cache: dict[str, Path],
    scale: float = 1.0,
) -> None:
    target_aspect = width / height
    cache_key = f"{image_path}|{target_aspect:.5f}"
    cropped_path = cache.get(cache_key)
    if cropped_path is None:
        cropped_path = crop_dir / f"{len(cache):05d}.png"
        center_crop_image(image_path, cropped_path, target_aspect)
        cache[cache_key] = cropped_path
    slide.shapes.add_picture(
        str(cropped_path),
        Inches(left * scale),
        Inches(top * scale),
        width=Inches(width * scale),
        height=Inches(height * scale),
    )


def add_image_contained(
    slide,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    background: RGBColor | None = None,
    scale: float = 1.0,
) -> None:
    if background is not None:
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left * scale),
            Inches(top * scale),
            Inches(width * scale),
            Inches(height * scale),
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = background
        cell.line.fill.background()
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    aspect = image_width / image_height
    cell_aspect = width / height
    if aspect >= cell_aspect:
        draw_width = width
        draw_height = width / aspect
    else:
        draw_height = height
        draw_width = height * aspect
    draw_left = left + (width - draw_width) / 2
    draw_top = top + (height - draw_height) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(draw_left * scale),
        Inches(draw_top * scale),
        width=Inches(draw_width * scale),
        height=Inches(draw_height * scale),
    )


def add_missing_cell(slide, left: float, top: float, width: float, height: float, text_color: RGBColor, scale: float = 1.0) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left * scale),
        Inches(top * scale),
        Inches(width * scale),
        Inches(height * scale),
    )
    shape.fill.background()
    shape.line.color.rgb = RGBColor(190, 190, 190)
    add_textbox(slide, "missing", left, top + height / 2 - 0.12, width, 0.24, 9, text_color, align=PP_ALIGN.CENTER, scale=scale)


def assignment_entries(settings: dict[str, Any]) -> list[dict[str, Any]]:
    raw = settings.get("assignments")
    if isinstance(raw, list):
        return [dict(item) for item in raw if isinstance(item, dict)]
    if isinstance(raw, dict):
        entries = []
        for image_id, value in raw.items():
            if isinstance(value, dict):
                entries.append({"image_id": image_id, **value})
        return entries
    return []


def resolve_assigned_groups(
    manifest: dict[str, Any],
    settings: dict[str, Any],
) -> tuple[list[str], list[dict[str, Any]], list[dict[str, str]]]:
    entries = assignment_entries(settings)
    if not entries:
        return [], [], []

    images = ordered_manifest_images(manifest, settings)
    by_id = {image["image_id"]: image for image in images}
    by_path = {str(image.get("relative_path") or ""): image for image in images}
    filename_counts: dict[str, int] = {}
    for image in images:
        filename = str(image.get("filename") or "")
        filename_counts[filename] = filename_counts.get(filename, 0) + 1
    for image in images:
        filename = str(image.get("filename") or "")
        if filename_counts.get(filename) == 1:
            by_path.setdefault(filename, image)

    order_rank = {image["image_id"]: index for index, image in enumerate(images)}
    resolved: list[tuple[int, dict[str, Any], str, str]] = []
    issues: list[dict[str, str]] = []
    used_images: set[str] = set()
    used_slots: set[tuple[str, str]] = set()
    for entry_index, entry in enumerate(entries):
        reference = str(entry.get("image_id") or entry.get("relative_path") or "").strip()
        image = by_id.get(reference) or by_path.get(reference)
        group_key = str(entry.get("group_key") or entry.get("group") or "").strip()
        channel_key = str(entry.get("channel_key") or entry.get("channel") or "").strip()
        if image is None:
            issues.append({"code": "unknown_image", "message": f"图片分配引用不存在：{reference or '(空)'}"})
            continue
        if not group_key or not channel_key:
            issues.append({"code": "incomplete_assignment", "message": f"{image['filename']} 缺少组别或通道"})
            continue
        image_id = image["image_id"]
        if image_id in used_images:
            issues.append({"code": "duplicate_image_assignment", "message": f"{image['filename']} 被重复分配"})
            continue
        slot = (group_key, channel_key)
        if slot in used_slots:
            issues.append({"code": "duplicate_assignment_slot", "message": f"组别 {group_key} 的 {channel_key} 位置被重复占用"})
            continue
        used_images.add(image_id)
        used_slots.add(slot)
        resolved.append((order_rank.get(image_id, len(images) + entry_index), image, group_key, channel_key))

    resolved.sort(key=lambda item: item[0])
    groups_by_key: dict[str, dict[str, Any]] = {}
    requested_groups = settings.get("group_order")
    group_order = []
    if isinstance(requested_groups, list):
        for value in requested_groups:
            key = str(value).strip()
            if key and key not in group_order:
                group_order.append(key)
    requested_channels = settings.get("channel_order")
    channel_order = []
    if isinstance(requested_channels, list):
        for value in requested_channels:
            key = str(value).strip()
            if key and key not in channel_order:
                channel_order.append(key)
    known_groups = {group["key"]: group.get("display") or group["key"] for group in manifest.get("groups", [])}
    for _, source, group_key, channel_key in resolved:
        group = groups_by_key.setdefault(
            group_key,
            {
                "key": group_key,
                "display": known_groups.get(group_key) or f"Group {(group_order.index(group_key) + 1) if group_key in group_order else len(groups_by_key) + 1}",
                "images": [],
            },
        )
        image = dict(source)
        image["channel"] = channel_key
        image["channel_label"] = channel_key
        image["channel_order"] = channel_order.index(channel_key) if channel_key in channel_order else len(channel_order)
        group["images"].append(image)
        if channel_key not in channel_order:
            channel_order.append(channel_key)
    ordered_group_keys = [key for key in group_order if key in groups_by_key]
    ordered_group_keys.extend(key for key in groups_by_key if key not in ordered_group_keys)
    for group in groups_by_key.values():
        group["images"].sort(key=lambda image: channel_order.index(image["channel"]) if image["channel"] in channel_order else len(channel_order))
    return channel_order, [groups_by_key[key] for key in ordered_group_keys], issues


def build_manual_groups(
    manifest: dict[str, Any],
    group_count: int,
    images_per_group: int,
    settings: dict[str, Any] | None = None,
) -> tuple[list[str], list[dict[str, Any]]]:
    group_count = min(max(group_count, 1), 50)
    images_per_group = min(max(images_per_group, 1), 20)
    channel_order = [f"slot{index + 1:02d}" for index in range(images_per_group)]
    flat_images = ordered_manifest_images(manifest, settings)
    groups: list[dict[str, Any]] = []

    for group_index in range(group_count):
        group_key = f"manual_group_{group_index + 1:02d}"
        group_images = []
        for image_index in range(images_per_group):
            flat_index = group_index * images_per_group + image_index
            if flat_index >= len(flat_images):
                continue
            image = dict(flat_images[flat_index])
            image["channel"] = channel_order[image_index]
            image["channel_label"] = channel_order[image_index]
            image["channel_order"] = image_index
            group_images.append(image)
        groups.append(
            {
                "key": group_key,
                "display": f"Group {group_index + 1}",
                "images": group_images,
            }
        )

    return channel_order, groups


def resolve_export_content(manifest: dict[str, Any], settings: dict[str, Any]) -> dict[str, Any]:
    layout_mode = settings.get("layout_mode") or "auto"
    rows_per_slide = int(settings.get("rows_per_slide") or 3)
    rows_per_slide = min(max(rows_per_slide, 1), 6)

    assigned_channels, assigned_groups, assignment_problems = resolve_assigned_groups(manifest, settings)
    if assignment_problems:
        raise ValueError(assignment_problems[0]["message"])

    if assigned_groups:
        requested_channels = settings.get("channel_order")
        channel_order = (
            [str(channel) for channel in requested_channels if str(channel)]
            if isinstance(requested_channels, list)
            else assigned_channels
        )
        selected = set(channel_order)
        groups = [
            {**group, "images": [image for image in group["images"] if image["channel"] in selected]}
            for group in assigned_groups
        ]
        rows_per_slide = min(rows_per_slide, max(len(groups), 1))
    elif layout_mode == "manual":
        group_count = int(settings.get("group_count") or 1)
        images_per_group = int(settings.get("images_per_group") or 1)
        manual_channel_order, groups = build_manual_groups(manifest, group_count, images_per_group, settings)
        requested_channels = settings.get("channel_order")
        channel_order = (
            [channel for channel in requested_channels if channel in manual_channel_order]
            if isinstance(requested_channels, list)
            else manual_channel_order
        )
        selected = set(channel_order)
        groups = [
            {**group, "images": [image for image in group["images"] if image["channel"] in selected]}
            for group in groups
        ]
        rows_per_slide = min(rows_per_slide, max(len(groups), 1))
    else:
        channel_order = settings.get("channel_order") or [channel["key"] for channel in manifest["channels"]]
        channel_order = [channel for channel in channel_order if channel]
        groups = manifest["groups"]

    if not channel_order:
        raise ValueError("预览和导出至少需要保留一个通道")

    return {
        "channel_order": channel_order,
        "groups": groups,
        "rows_per_slide": rows_per_slide,
        "labels": settings.get("labels") or {},
        "group_labels": settings.get("group_labels") or {},
        "background": settings.get("background") or "white",
        "title": (settings.get("title") or "").strip(),
        "panel_letter": (settings.get("panel_letter") or "").strip(),
        "show_group_labels": bool(settings.get("show_sample_name", True)),
        "group_label_side": settings.get("group_label_side") or "left",
        "fit_mode": settings.get("fit_mode") or "crop",
    }


def calculate_layout(
    channel_order: list[str],
    rows_per_slide: int,
    title: str,
    show_group_labels: bool,
    group_label_side: str,
    slide_width: float = 13.333,
    slide_height: float = 7.5,
    coordinate_scale: float = 1.0,
) -> dict[str, float]:
    margin_left = 0.45
    margin_right = 0.4
    margin_bottom = 0.35
    top_margin = 0.22
    title_text_height = max(0.28, 0.12 / coordinate_scale)
    marker_text_height = max(0.24, 0.12 / coordinate_scale)
    title_height = title_text_height + 0.06 / coordinate_scale if title else 0
    marker_height = marker_text_height + 0.08 / coordinate_scale
    side_label_width = 0.42 if show_group_labels else 0
    side_gap = 0.08 if show_group_labels else 0
    col_gap = 0.04
    row_gap = 0.08
    col_count = len(channel_order)

    if group_label_side == "right":
        grid_left = margin_left
        grid_right = slide_width - margin_right - side_label_width - side_gap
    else:
        grid_left = margin_left + side_label_width + side_gap
        grid_right = slide_width - margin_right

    grid_available_width = grid_right - grid_left
    grid_top = top_margin + title_height + marker_height
    grid_available_height = slide_height - grid_top - margin_bottom
    cell_by_width = (grid_available_width - col_gap * (col_count - 1)) / col_count
    cell_by_height = (grid_available_height - row_gap * (rows_per_slide - 1)) / rows_per_slide
    cell_size = max(0.55, min(cell_by_width, cell_by_height))
    actual_grid_width = cell_size * col_count + col_gap * (col_count - 1)
    actual_grid_left = grid_left + max(0, (grid_available_width - actual_grid_width) / 2)

    return {
        "slide_width": slide_width,
        "slide_height": slide_height,
        "margin_left": margin_left,
        "margin_right": margin_right,
        "top_margin": top_margin,
        "title_height": title_height,
        "title_text_height": title_text_height,
        "marker_height": marker_height,
        "marker_text_height": marker_text_height,
        "side_label_width": side_label_width,
        "side_gap": side_gap,
        "col_gap": col_gap,
        "row_gap": row_gap,
        "grid_top": grid_top,
        "cell_size": cell_size,
        "actual_grid_width": actual_grid_width,
        "actual_grid_left": actual_grid_left,
    }


def build_pptx(manifest: dict[str, Any], settings: dict[str, Any]) -> Path:
    content = resolve_export_content(manifest, settings)
    profile = resolve_export_profile(settings)
    coordinate_scale = profile["coordinate_scale"]
    channel_order = content["channel_order"]
    groups = content["groups"]
    labels = content["labels"]
    group_labels = content["group_labels"]
    rows_per_slide = content["rows_per_slide"]
    background = content["background"]
    title = content["title"]
    panel_letter = content["panel_letter"]
    show_group_labels = content["show_group_labels"]
    group_label_side = content["group_label_side"]
    fit_mode = content["fit_mode"]

    dark = background == "black"
    bg_color = RGBColor(0, 0, 0) if dark else RGBColor(255, 255, 255)
    text_color = RGBColor(245, 245, 245) if dark else RGBColor(15, 15, 15)
    secondary_color = RGBColor(210, 210, 210) if dark else RGBColor(45, 45, 45)

    prs = Presentation()
    prs.slide_width = Inches(profile["width_in"])
    prs.slide_height = Inches(profile["height_in"])
    blank_layout = prs.slide_layouts[6]

    layout = calculate_layout(
        channel_order,
        rows_per_slide,
        title,
        show_group_labels,
        group_label_side,
        slide_height=profile["logical_height"],
        coordinate_scale=coordinate_scale,
    )
    slide_width = layout["slide_width"]
    margin_left = layout["margin_left"]
    margin_right = layout["margin_right"]
    top_margin = layout["top_margin"]
    title_height = layout["title_height"]
    side_label_width = layout["side_label_width"]
    side_gap = layout["side_gap"]
    col_gap = layout["col_gap"]
    row_gap = layout["row_gap"]
    grid_top = layout["grid_top"]
    cell_size = layout["cell_size"]
    actual_grid_width = layout["actual_grid_width"]
    actual_grid_left = layout["actual_grid_left"]

    output_name = f"fluorescence_batch_{unique_export_stamp()}.pptx"
    output_path = job_export_dir(manifest) / output_name
    crop_dir = JOBS_DIR / manifest["job_id"] / "export_crops" / uuid.uuid4().hex[:8]
    crop_cache: dict[str, Path] = {}

    for page_start in range(0, len(groups), rows_per_slide):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        title_left = margin_left + (0.36 if panel_letter else 0)
        if panel_letter:
            add_textbox(slide, panel_letter, margin_left, top_margin, 0.28, layout["title_text_height"], 16, text_color, bold=True, scale=coordinate_scale)
        if title:
            add_textbox(slide, title, title_left, top_margin, slide_width - title_left - margin_right, layout["title_text_height"], 14, text_color, bold=True, scale=coordinate_scale)

        marker_top = top_margin + title_height
        for column_index, channel in enumerate(channel_order):
            col_left = actual_grid_left + column_index * (cell_size + col_gap)
            label = labels.get(channel) or channel
            add_textbox(
                slide,
                label,
                col_left,
                marker_top,
                cell_size,
                layout["marker_text_height"],
                12,
                text_color,
                bold=True,
                align=PP_ALIGN.CENTER,
                scale=coordinate_scale,
            )

        page_groups = groups[page_start : page_start + rows_per_slide]
        for row_index, group in enumerate(page_groups):
            row_top = grid_top + row_index * (cell_size + row_gap)
            if show_group_labels:
                group_label = group_labels.get(group["key"]) or group["display"]
                if group_label_side == "right":
                    label_left = actual_grid_left + actual_grid_width + side_gap
                    rotation = 90
                else:
                    label_left = margin_left
                    rotation = 270
                add_textbox(
                    slide,
                    group_label,
                    label_left,
                    row_top,
                    side_label_width,
                    cell_size,
                    12,
                    secondary_color,
                    bold=True,
                    align=PP_ALIGN.CENTER,
                    rotation=rotation,
                    scale=coordinate_scale,
                )

            images_by_channel = {image["channel"]: image for image in group["images"]}
            for column_index, channel in enumerate(channel_order):
                col_left = actual_grid_left + column_index * (cell_size + col_gap)
                image = images_by_channel.get(channel)
                if image:
                    source_path = image_source_path(image)
                    if fit_mode == "contain":
                        add_image_contained(slide, source_path, col_left, row_top, cell_size, cell_size, RGBColor(0, 0, 0), coordinate_scale)
                    else:
                        add_image_cover(
                            slide,
                            source_path,
                            col_left,
                            row_top,
                            cell_size,
                            cell_size,
                            crop_dir,
                            crop_cache,
                            coordinate_scale,
                        )
                    add_scale_bar_pptx(
                        slide,
                        source_path,
                        (col_left, row_top, cell_size, cell_size),
                        fit_mode,
                        settings,
                        coordinate_scale,
                        image,
                    )
                else:
                    add_missing_cell(slide, col_left, row_top, cell_size, cell_size, secondary_color, coordinate_scale)

    prs.save(output_path)
    return output_path


def add_ppt_picture_cover(
    slide,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    crop_dir: Path,
    cache: dict[str, Path],
    scale: float = 1.0,
) -> None:
    add_image_cover(slide, image_path, left, top, width, height, crop_dir, cache, scale)


def add_ppt_line(
    slide,
    start: tuple[float, float],
    end: tuple[float, float],
    color: RGBColor = RGBColor(60, 60, 60),
    scale: float = 1.0,
    width_pt: float = 1.0,
) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start[0] * scale),
        Inches(start[1] * scale),
        Inches(end[0] * scale),
        Inches(end[1] * scale),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)


def add_ppt_roi(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    color: RGBColor = RGBColor(45, 45, 45),
    scale: float = 1.0,
) -> None:
    roi = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left * scale),
        Inches(top * scale),
        Inches(width * scale),
        Inches(height * scale),
    )
    roi.fill.background()
    roi.line.color.rgb = color
    roi.line.width = Pt(1.2)


def build_ihc_groups(
    manifest: dict[str, Any],
    group_count: int,
    settings: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    group_count = min(max(group_count, 1), 50)
    if settings:
        _, assigned_groups, problems = resolve_assigned_groups(manifest, settings)
        if problems:
            raise ValueError(problems[0]["message"])
        if assigned_groups:
            unknown = sorted(
                {image["channel"] for group in assigned_groups for image in group["images"]} - {"low", "high"}
            )
            if unknown:
                raise ValueError(f"IHC 分配的通道只能是 low 或 high：{', '.join(unknown)}")
            return assigned_groups

    flat_images = ordered_manifest_images(manifest, settings)
    groups: list[dict[str, Any]] = []
    for group_index in range(group_count):
        group_images = []
        for image_index, channel in enumerate(["low", "high"]):
            flat_index = group_index * 2 + image_index
            if flat_index >= len(flat_images):
                continue
            image = dict(flat_images[flat_index])
            image["channel"] = channel
            image["channel_label"] = channel
            image["channel_order"] = image_index
            group_images.append(image)
        groups.append(
            {
                "key": f"ihc_group_{group_index + 1:02d}",
                "display": f"Group {group_index + 1}",
                "images": group_images,
            }
        )
    return groups


def normalized_ihc_roi(settings: dict[str, Any], group_key: str | None = None) -> tuple[float, float, float, float]:
    roi_settings = settings
    raw_rois = settings.get("ihc_rois")
    if group_key and isinstance(raw_rois, list):
        match = next(
            (item for item in raw_rois if isinstance(item, dict) and str(item.get("group_key") or "") == group_key),
            None,
        )
        if match:
            roi_settings = {
                "ihc_roi_x": match.get("x", match.get("ihc_roi_x")),
                "ihc_roi_y": match.get("y", match.get("ihc_roi_y")),
                "ihc_roi_w": match.get("w", match.get("ihc_roi_w")),
                "ihc_roi_h": match.get("h", match.get("ihc_roi_h")),
            }
    elif group_key and isinstance(raw_rois, dict):
        match = raw_rois.get(group_key)
        if isinstance(match, dict):
            roi_settings = {
                "ihc_roi_x": match.get("x", match.get("ihc_roi_x")),
                "ihc_roi_y": match.get("y", match.get("ihc_roi_y")),
                "ihc_roi_w": match.get("w", match.get("ihc_roi_w")),
                "ihc_roi_h": match.get("h", match.get("ihc_roi_h")),
            }

    def value(key: str, default: float) -> float:
        raw = roi_settings.get(key)
        return float(default if raw is None or raw == "" else raw)

    width = min(max(value("ihc_roi_w", 0.30), 0.05), 1.0)
    height = min(max(value("ihc_roi_h", 0.28), 0.05), 1.0)
    x = min(max(value("ihc_roi_x", 0.32), 0.0), 1.0 - width)
    y = min(max(value("ihc_roi_y", 0.36), 0.0), 1.0 - height)
    return x, y, width, height


def build_ihc_pptx(manifest: dict[str, Any], settings: dict[str, Any]) -> Path:
    profile = resolve_export_profile(settings)
    coordinate_scale = profile["coordinate_scale"]
    group_count = int(settings.get("group_count") or 1)
    groups = build_ihc_groups(manifest, group_count, settings)
    group_labels = settings.get("group_labels") or {}
    low_label = (settings.get("ihc_low_label") or "4X").strip()
    high_label = (settings.get("ihc_high_label") or "20X").strip()
    title = (settings.get("title") or "").strip()
    panel_letter = (settings.get("panel_letter") or "").strip()
    draw_connectors = bool(settings.get("ihc_draw_connectors", True))
    rows_per_slide = min(max(int(settings.get("rows_per_slide") or len(groups) or 1), 1), 4)
    dark = (settings.get("background") or "white") == "black"
    show_group_labels = bool(settings.get("show_sample_name", True))
    group_label_side = settings.get("group_label_side") or "left"
    fit_mode = settings.get("fit_mode") or "crop"
    bg_color = RGBColor(0, 0, 0) if dark else RGBColor(255, 255, 255)
    text_color = RGBColor(245, 245, 245) if dark else RGBColor(0, 0, 0)
    secondary_color = RGBColor(210, 210, 210) if dark else RGBColor(80, 80, 80)
    line_color = RGBColor(225, 225, 225) if dark else RGBColor(60, 60, 60)

    prs = Presentation()
    prs.slide_width = Inches(profile["width_in"])
    prs.slide_height = Inches(profile["height_in"])
    blank_layout = prs.slide_layouts[6]

    slide_width = 13.333
    slide_height = profile["logical_height"]
    margin_left = 0.45
    margin_right = 0.45
    top_margin = 0.24
    title_text_height = max(0.28, 0.12 / coordinate_scale)
    header_text_height = max(0.28, 0.12 / coordinate_scale)
    title_height = title_text_height + 0.06 / coordinate_scale if title else 0
    header_height = header_text_height + 0.08 / coordinate_scale
    bottom_margin = 0.35
    row_gap = 0.34
    side_label_width = 0.55 if show_group_labels else 0
    side_gap = 0.18 if show_group_labels else 0
    between_cols = 0.85
    content_top = top_margin + title_height + header_height
    row_height = (slide_height - content_top - bottom_margin - row_gap * (rows_per_slide - 1)) / rows_per_slide
    image_height = max(0.72, row_height)
    image_width = image_height * 1.72
    low_left = margin_left + (side_label_width + side_gap if show_group_labels and group_label_side != "right" else 0)
    high_left = low_left + image_width + between_cols
    right_reserved = side_label_width + side_gap if show_group_labels and group_label_side == "right" else 0
    if high_left + image_width > slide_width - margin_right - right_reserved:
        available = slide_width - margin_right - right_reserved - low_left - between_cols
        image_width = available / 2
        image_height = min(image_height, image_width / 1.72)
        high_left = low_left + image_width + between_cols

    output_name = f"ihc_batch_{unique_export_stamp()}.pptx"
    output_path = job_export_dir(manifest) / output_name
    crop_dir = JOBS_DIR / manifest["job_id"] / "ihc_export_crops" / uuid.uuid4().hex[:8]
    crop_cache: dict[str, Path] = {}

    for page_start in range(0, len(groups), rows_per_slide):
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        title_left = margin_left + (0.34 if panel_letter else 0)
        if panel_letter:
            add_textbox(slide, panel_letter, margin_left, top_margin, 0.28, title_text_height, 16, text_color, bold=True, scale=coordinate_scale)
        if title:
            add_textbox(slide, title, title_left, top_margin, slide_width - title_left - margin_right, title_text_height, 14, text_color, bold=True, scale=coordinate_scale)

        header_top = top_margin + title_height
        add_textbox(slide, low_label, low_left, header_top, image_width, header_text_height, 14, text_color, bold=True, align=PP_ALIGN.CENTER, scale=coordinate_scale)
        add_textbox(slide, high_label, high_left, header_top, image_width, header_text_height, 14, text_color, bold=True, align=PP_ALIGN.CENTER, scale=coordinate_scale)

        page_groups = groups[page_start : page_start + rows_per_slide]
        for row_index, group in enumerate(page_groups):
            roi_x, roi_y, roi_w, roi_h = normalized_ihc_roi(settings, group["key"])
            row_top = content_top + row_index * (row_height + row_gap)
            row_top += max(0, (row_height - image_height) / 2)
            label = group_labels.get(group["key"]) or group["display"]
            if show_group_labels:
                label_left = high_left + image_width + side_gap if group_label_side == "right" else margin_left
                rotation = 90 if group_label_side == "right" else 270
                add_textbox(slide, label, label_left, row_top, side_label_width, image_height, 13, text_color, bold=True, align=PP_ALIGN.CENTER, rotation=rotation, scale=coordinate_scale)
            images_by_channel = {image["channel"]: image for image in group["images"]}
            low_image = images_by_channel.get("low")
            high_image = images_by_channel.get("high")
            if low_image:
                low_source = image_source_path(low_image)
                if fit_mode == "contain":
                    add_image_contained(slide, low_source, low_left, row_top, image_width, image_height, bg_color, coordinate_scale)
                else:
                    add_ppt_picture_cover(slide, low_source, low_left, row_top, image_width, image_height, crop_dir, crop_cache, coordinate_scale)
                add_scale_bar_pptx(slide, low_source, (low_left, row_top, image_width, image_height), fit_mode, settings, coordinate_scale, low_image)
            else:
                add_missing_cell(slide, low_left, row_top, image_width, image_height, secondary_color, coordinate_scale)
            if high_image:
                high_source = image_source_path(high_image)
                if fit_mode == "contain":
                    add_image_contained(slide, high_source, high_left, row_top, image_width, image_height, bg_color, coordinate_scale)
                else:
                    add_ppt_picture_cover(slide, high_source, high_left, row_top, image_width, image_height, crop_dir, crop_cache, coordinate_scale)
                add_scale_bar_pptx(slide, high_source, (high_left, row_top, image_width, image_height), fit_mode, settings, coordinate_scale, high_image)
            else:
                add_missing_cell(slide, high_left, row_top, image_width, image_height, secondary_color, coordinate_scale)

            roi_left = low_left + roi_x * image_width
            roi_top = row_top + roi_y * image_height
            roi_width = roi_w * image_width
            roi_height = roi_h * image_height
            add_ppt_roi(slide, roi_left, roi_top, roi_width, roi_height, line_color, coordinate_scale)
            if draw_connectors:
                add_ppt_line(slide, (roi_left + roi_width, roi_top), (high_left, row_top), line_color, coordinate_scale)
                add_ppt_line(slide, (roi_left + roi_width, roi_top + roi_height), (high_left, row_top + image_height), line_color, coordinate_scale)

    prs.save(output_path)
    return output_path


def font(size_px: int, bold: bool = False) -> ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/msttcorefonts/Arial_Bold.ttf" if bold else "/usr/share/fonts/truetype/msttcorefonts/Arial.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size_px)
    return ImageFont.load_default()


def inch(value: float) -> int:
    return int(round(value * RASTER_DPI))


def rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return (color[0], color[1], color[2])


def draw_centered_text(
    canvas: Image.Image,
    text: str,
    box: tuple[int, int, int, int],
    fill: tuple[int, int, int],
    size_px: int,
    bold: bool = False,
    rotation: int | None = None,
) -> None:
    if not text:
        return
    draw = ImageDraw.Draw(canvas)
    typeface = font(size_px, bold=bold)

    if rotation is None:
        bbox = draw.textbbox((0, 0), text, font=typeface)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        x = box[0] + (box[2] - box[0] - text_width) / 2
        y = box[1] + (box[3] - box[1] - text_height) / 2 - bbox[1]
        draw.text((x, y), text, fill=fill, font=typeface)
        return

    width = max(1, box[2] - box[0])
    height = max(1, box[3] - box[1])
    layer = Image.new("RGBA", (height, width), (255, 255, 255, 0))
    layer_draw = ImageDraw.Draw(layer)
    bbox = layer_draw.textbbox((0, 0), text, font=typeface)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    x = (height - text_width) / 2
    y = (width - text_height) / 2 - bbox[1]
    layer_draw.text((x, y), text, fill=(*fill, 255), font=typeface)
    rotated = layer.rotate(rotation, expand=True)
    canvas.paste(rotated, (box[0], box[1]), rotated)


def draw_left_text(
    canvas: Image.Image,
    text: str,
    box: tuple[int, int, int, int],
    fill: tuple[int, int, int],
    size_px: int,
    bold: bool = False,
) -> None:
    if not text:
        return
    draw = ImageDraw.Draw(canvas)
    typeface = font(size_px, bold=bold)
    bbox = draw.textbbox((0, 0), text, font=typeface)
    y = box[1] + (box[3] - box[1] - (bbox[3] - bbox[1])) / 2 - bbox[1]
    draw.text((box[0], y), text, fill=fill, font=typeface)


def resize_cover(image: Image.Image, size: tuple[int, int]) -> Image.Image:
    width, height = image.size
    target_width, target_height = size
    target_aspect = target_width / target_height
    source_aspect = width / height
    if source_aspect > target_aspect:
        new_width = int(height * target_aspect)
        left = max(0, (width - new_width) // 2)
        image = image.crop((left, 0, left + new_width, height))
    else:
        new_height = int(width / target_aspect)
        top = max(0, (height - new_height) // 2)
        image = image.crop((0, top, width, top + new_height))
    return image.resize(size, Image.Resampling.LANCZOS)


def resize_contain(
    image: Image.Image,
    size: tuple[int, int],
    background: tuple[int, int, int] = (0, 0, 0),
) -> Image.Image:
    target = Image.new("RGB", size, background)
    copy = image.copy()
    copy.thumbnail(size, Image.Resampling.LANCZOS)
    left = (size[0] - copy.size[0]) // 2
    top = (size[1] - copy.size[1]) // 2
    target.paste(copy, (left, top))
    return target


def build_raster_pages(manifest: dict[str, Any], settings: dict[str, Any]) -> list[Image.Image]:
    content = resolve_export_content(manifest, settings)
    profile = resolve_export_profile(settings)
    dpi = profile["dpi"]
    logical_pixel_scale = profile["pixel_width"] / profile["logical_width"]
    px = lambda value: max(0, int(round(value * logical_pixel_scale)))
    text_px = lambda value: max(1, int(round(max(7, value * 72 / RASTER_DPI * min(1, profile["coordinate_scale"])) * dpi / 72)))
    channel_order = content["channel_order"]
    groups = content["groups"]
    labels = content["labels"]
    group_labels = content["group_labels"]
    rows_per_slide = content["rows_per_slide"]
    title = content["title"]
    panel_letter = content["panel_letter"]
    show_group_labels = content["show_group_labels"]
    group_label_side = content["group_label_side"]
    fit_mode = content["fit_mode"]

    dark = content["background"] == "black"
    bg_color = (0, 0, 0) if dark else (255, 255, 255)
    text_color = (245, 245, 245) if dark else (15, 15, 15)
    secondary_color = (210, 210, 210) if dark else (45, 45, 45)

    layout = calculate_layout(
        channel_order,
        rows_per_slide,
        title,
        show_group_labels,
        group_label_side,
        slide_height=profile["logical_height"],
        coordinate_scale=profile["coordinate_scale"],
    )
    canvas_size = (profile["pixel_width"], profile["pixel_height"])
    pages: list[Image.Image] = []

    for page_start in range(0, len(groups), rows_per_slide):
        canvas = Image.new("RGB", canvas_size, bg_color)

        title_left = layout["margin_left"] + (0.36 if panel_letter else 0)
        if panel_letter:
            draw_left_text(
                canvas,
                panel_letter,
                (px(layout["margin_left"]), px(layout["top_margin"]), px(layout["margin_left"] + 0.28), px(layout["top_margin"] + layout["title_text_height"])),
                text_color,
                text_px(52),
                bold=True,
            )
        if title:
            draw_left_text(
                canvas,
                title,
                (px(title_left), px(layout["top_margin"]), px(layout["slide_width"] - layout["margin_right"]), px(layout["top_margin"] + layout["title_text_height"])),
                text_color,
                text_px(46),
                bold=True,
            )

        marker_top = layout["top_margin"] + layout["title_height"]
        for column_index, channel in enumerate(channel_order):
            col_left = layout["actual_grid_left"] + column_index * (layout["cell_size"] + layout["col_gap"])
            label = labels.get(channel) or channel
            draw_centered_text(
                canvas,
                label,
                (px(col_left), px(marker_top), px(col_left + layout["cell_size"]), px(marker_top + layout["marker_text_height"])),
                text_color,
                text_px(40),
                bold=True,
            )

        page_groups = groups[page_start : page_start + rows_per_slide]
        for row_index, group in enumerate(page_groups):
            row_top = layout["grid_top"] + row_index * (layout["cell_size"] + layout["row_gap"])
            if show_group_labels:
                group_label = group_labels.get(group["key"]) or group["display"]
                if group_label_side == "right":
                    label_left = layout["actual_grid_left"] + layout["actual_grid_width"] + layout["side_gap"]
                    rotation = 270
                else:
                    label_left = layout["margin_left"]
                    rotation = 90
                draw_centered_text(
                    canvas,
                    group_label,
                    (px(label_left), px(row_top), px(label_left + layout["side_label_width"]), px(row_top + layout["cell_size"])),
                    secondary_color,
                    text_px(40),
                    bold=True,
                    rotation=rotation,
                )

            images_by_channel = {image["channel"]: image for image in group["images"]}
            for column_index, channel in enumerate(channel_order):
                col_left = layout["actual_grid_left"] + column_index * (layout["cell_size"] + layout["col_gap"])
                box = (
                    px(col_left),
                    px(row_top),
                    px(col_left + layout["cell_size"]),
                    px(row_top + layout["cell_size"]),
                )
                image = images_by_channel.get(channel)
                if image:
                    with Image.open(image_source_path(image)) as source:
                        rgb = image_to_rgb(source)
                        if fit_mode == "contain":
                            rendered = resize_contain(rgb, (box[2] - box[0], box[3] - box[1]))
                        else:
                            rendered = resize_cover(rgb, (box[2] - box[0], box[3] - box[1]))
                    canvas.paste(rendered, (box[0], box[1]))
                    draw_scale_bar_raster(canvas, image_source_path(image), box, fit_mode, settings, dpi, image)
                else:
                    draw = ImageDraw.Draw(canvas)
                    draw.rectangle(box, outline=(190, 190, 190), width=2)
                    draw_centered_text(canvas, "missing", box, secondary_color, text_px(26))

        pages.append(canvas)

    return pages


def draw_rotated_label(
    canvas: Image.Image,
    text: str,
    box: tuple[int, int, int, int],
    fill: tuple[int, int, int] = (0, 0, 0),
    rotation: int = 90,
    size_px: int = 42,
) -> None:
    draw_centered_text(canvas, text, box, fill, size_px, bold=True, rotation=rotation)


def build_ihc_raster_pages(manifest: dict[str, Any], settings: dict[str, Any]) -> list[Image.Image]:
    profile = resolve_export_profile(settings)
    dpi = profile["dpi"]
    logical_pixel_scale = profile["pixel_width"] / profile["logical_width"]
    px = lambda value: max(0, int(round(value * logical_pixel_scale)))
    text_px = lambda value: max(1, int(round(max(7, value * 72 / RASTER_DPI * min(1, profile["coordinate_scale"])) * dpi / 72)))
    group_count = int(settings.get("group_count") or 1)
    groups = build_ihc_groups(manifest, group_count, settings)
    group_labels = settings.get("group_labels") or {}
    low_label = (settings.get("ihc_low_label") or "4X").strip()
    high_label = (settings.get("ihc_high_label") or "20X").strip()
    title = (settings.get("title") or "").strip()
    panel_letter = (settings.get("panel_letter") or "").strip()
    draw_connectors = bool(settings.get("ihc_draw_connectors", True))
    rows_per_slide = min(max(int(settings.get("rows_per_slide") or len(groups) or 1), 1), 4)
    dark = (settings.get("background") or "white") == "black"
    show_group_labels = bool(settings.get("show_sample_name", True))
    group_label_side = settings.get("group_label_side") or "left"
    fit_mode = settings.get("fit_mode") or "crop"
    bg_color = (0, 0, 0) if dark else (255, 255, 255)
    text_color = (245, 245, 245) if dark else (0, 0, 0)
    secondary_color = (210, 210, 210) if dark else (80, 80, 80)
    line_color = (225, 225, 225) if dark else (60, 60, 60)

    slide_width = 13.333
    slide_height = profile["logical_height"]
    margin_left = 0.45
    margin_right = 0.45
    top_margin = 0.24
    coordinate_scale = profile["coordinate_scale"]
    title_text_height = max(0.28, 0.12 / coordinate_scale)
    header_text_height = max(0.28, 0.12 / coordinate_scale)
    title_height = title_text_height + 0.06 / coordinate_scale if title else 0
    header_height = header_text_height + 0.08 / coordinate_scale
    bottom_margin = 0.35
    row_gap = 0.34
    side_label_width = 0.55 if show_group_labels else 0
    side_gap = 0.18 if show_group_labels else 0
    between_cols = 0.85
    content_top = top_margin + title_height + header_height
    row_height = (slide_height - content_top - bottom_margin - row_gap * (rows_per_slide - 1)) / rows_per_slide
    image_height = max(0.72, row_height)
    image_width = image_height * 1.72
    low_left = margin_left + (side_label_width + side_gap if show_group_labels and group_label_side != "right" else 0)
    high_left = low_left + image_width + between_cols
    right_reserved = side_label_width + side_gap if show_group_labels and group_label_side == "right" else 0
    if high_left + image_width > slide_width - margin_right - right_reserved:
        available = slide_width - margin_right - right_reserved - low_left - between_cols
        image_width = available / 2
        image_height = min(image_height, image_width / 1.72)
        high_left = low_left + image_width + between_cols

    canvas_size = (profile["pixel_width"], profile["pixel_height"])
    pages: list[Image.Image] = []

    for page_start in range(0, len(groups), rows_per_slide):
        canvas = Image.new("RGB", canvas_size, bg_color)
        draw = ImageDraw.Draw(canvas)

        title_left = margin_left + (0.34 if panel_letter else 0)
        if panel_letter:
            draw_left_text(canvas, panel_letter, (px(margin_left), px(top_margin), px(margin_left + 0.28), px(top_margin + title_text_height)), text_color, text_px(52), bold=True)
        if title:
            draw_left_text(canvas, title, (px(title_left), px(top_margin), px(slide_width - margin_right), px(top_margin + title_text_height)), text_color, text_px(46), bold=True)

        header_top = top_margin + title_height
        draw_centered_text(canvas, low_label, (px(low_left), px(header_top), px(low_left + image_width), px(header_top + header_text_height)), text_color, text_px(44), bold=True)
        draw_centered_text(canvas, high_label, (px(high_left), px(header_top), px(high_left + image_width), px(header_top + header_text_height)), text_color, text_px(44), bold=True)

        page_groups = groups[page_start : page_start + rows_per_slide]
        for row_index, group in enumerate(page_groups):
            roi_x, roi_y, roi_w, roi_h = normalized_ihc_roi(settings, group["key"])
            row_top = content_top + row_index * (row_height + row_gap)
            row_top += max(0, (row_height - image_height) / 2)
            label = group_labels.get(group["key"]) or group["display"]
            if show_group_labels:
                label_left = high_left + image_width + side_gap if group_label_side == "right" else margin_left
                rotation = 270 if group_label_side == "right" else 90
                draw_rotated_label(
                    canvas,
                    label,
                    (px(label_left), px(row_top), px(label_left + side_label_width), px(row_top + image_height)),
                    text_color,
                    rotation,
                    text_px(42),
                )
            images_by_channel = {image["channel"]: image for image in group["images"]}
            boxes = {
                "low": (px(low_left), px(row_top), px(low_left + image_width), px(row_top + image_height)),
                "high": (px(high_left), px(row_top), px(high_left + image_width), px(row_top + image_height)),
            }
            for channel in ["low", "high"]:
                image = images_by_channel.get(channel)
                box = boxes[channel]
                if image:
                    with Image.open(image_source_path(image)) as source:
                        rgb = image_to_rgb(source)
                        if fit_mode == "contain":
                            rendered = resize_contain(rgb, (box[2] - box[0], box[3] - box[1]), bg_color)
                        else:
                            rendered = resize_cover(rgb, (box[2] - box[0], box[3] - box[1]))
                    canvas.paste(rendered, (box[0], box[1]))
                    draw_scale_bar_raster(canvas, image_source_path(image), box, fit_mode, settings, dpi, image)
                else:
                    draw.rectangle(box, outline=secondary_color, width=2)
                    draw_centered_text(canvas, "missing", box, secondary_color, text_px(26))

            low_box = boxes["low"]
            high_box = boxes["high"]
            roi_box = (
                int(low_box[0] + roi_x * (low_box[2] - low_box[0])),
                int(low_box[1] + roi_y * (low_box[3] - low_box[1])),
                int(low_box[0] + (roi_x + roi_w) * (low_box[2] - low_box[0])),
                int(low_box[1] + (roi_y + roi_h) * (low_box[3] - low_box[1])),
            )
            line_width = max(2, round(3 * dpi / RASTER_DPI))
            draw.rectangle(roi_box, outline=line_color, width=line_width)
            if draw_connectors:
                draw.line((roi_box[2], roi_box[1], high_box[0], high_box[1]), fill=line_color, width=line_width)
                draw.line((roi_box[2], roi_box[3], high_box[0], high_box[3]), fill=line_color, width=line_width)

        pages.append(canvas)

    return pages


def build_raster_export(manifest: dict[str, Any], settings: dict[str, Any], export_format: str) -> Path:
    figure_type = (settings.get("figure_type") or "if").lower()
    profile = resolve_export_profile(settings)
    dpi = profile["dpi"]
    pages = build_ihc_raster_pages(manifest, settings) if figure_type == "ihc" else build_raster_pages(manifest, settings)
    if not pages:
        raise ValueError("没有可导出的页面")

    output_dir = job_export_dir(manifest)
    timestamp = unique_export_stamp()
    prefix = "ihc_batch" if figure_type == "ihc" else "fluorescence_batch"

    if export_format == "pdf":
        output_path = output_dir / f"{prefix}_{timestamp}.pdf"
        pages[0].save(output_path, "PDF", save_all=True, append_images=pages[1:], resolution=dpi)
        return output_path

    if export_format not in {"png", "jpg", "tif", "tiff"}:
        raise ValueError(f"Unsupported export format: {export_format}")

    extension = "tif" if export_format in {"tif", "tiff"} else export_format
    image_format = {"jpg": "JPEG", "png": "PNG", "tif": "TIFF"}[extension]

    def save_page(page: Image.Image, path: Path) -> None:
        if image_format == "JPEG":
            page.save(path, image_format, quality=95, dpi=(dpi, dpi))
        elif image_format == "TIFF":
            page.save(path, image_format, compression="raw", dpi=(dpi, dpi))
        else:
            page.save(path, image_format, dpi=(dpi, dpi))

    if len(pages) == 1:
        output_path = output_dir / f"{prefix}_{timestamp}.{extension}"
        save_page(pages[0], output_path)
        return output_path

    zip_path = output_dir / f"{prefix}_{timestamp}_{extension}.zip"
    page_dir = output_dir / f"{prefix}_{timestamp}_{extension}_pages"
    page_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as archive:
        for index, page in enumerate(pages, start=1):
            page_path = page_dir / f"page_{index:02d}.{extension}"
            save_page(page, page_path)
            archive.write(page_path, page_path.name)
    return zip_path


def preflight_export(manifest: dict[str, Any], settings: dict[str, Any]) -> dict[str, Any]:
    errors: list[dict[str, str]] = []
    warnings: list[dict[str, str]] = []

    def add(target: list[dict[str, str]], code: str, message: str) -> None:
        target.append({"code": code, "message": message})

    try:
        profile = resolve_export_profile(settings)
    except ValueError as exc:
        profile = None
        add(errors, "invalid_export_profile", str(exc))

    figure_type = str(settings.get("figure_type") or "if").lower()
    if figure_type not in {"if", "ihc"}:
        add(errors, "invalid_figure_type", "图版类型仅支持 IF 或 IHC")
    export_format = str(settings.get("export_format") or "pptx").lower()
    if export_format not in {"pptx", "png", "jpg", "pdf", "tif", "tiff", "zip"}:
        add(errors, "invalid_export_format", "不支持所选导出格式")

    entries = assignment_entries(settings)
    assigned_channels, assigned_groups, assignment_problems = resolve_assigned_groups(manifest, settings)
    errors.extend(assignment_problems)
    if entries and not assigned_groups:
        add(errors, "no_valid_assignments", "没有可用的图片分配")
    if assigned_groups:
        assigned_ids = {image["image_id"] for group in assigned_groups for image in group["images"]}
        if len(assigned_ids) < len(manifest.get("images", [])):
            add(warnings, "unassigned_images", f"有 {len(manifest.get('images', [])) - len(assigned_ids)} 张原图未放入图版")
    if entries and assigned_groups:
        requested_group_order = settings.get("group_order")
        group_order = [str(value).strip() for value in requested_group_order if str(value).strip()] if isinstance(requested_group_order, list) else []
        if len(group_order) != len(set(group_order)):
            add(errors, "duplicate_group_order", "组别顺序中存在重复项")
        expected_group_count = min(max(int(settings.get("group_count") or len(group_order) or len(assigned_groups)), 1), 50)
        assigned_group_keys = {group["key"] for group in assigned_groups}
        if group_order:
            missing_groups = [key for key in group_order if key not in assigned_group_keys]
            extra_groups = sorted(assigned_group_keys - set(group_order))
            for key in missing_groups:
                add(errors, "missing_assignment_group", f"组别 {key} 没有任何已映射图片")
            for key in extra_groups:
                add(errors, "unknown_assignment_group", f"图片被映射到未声明组别 {key}")
            groups_to_check = group_order
        else:
            groups_to_check = [group["key"] for group in assigned_groups]
            if len(assigned_group_keys) != expected_group_count:
                add(errors, "incomplete_group_count", f"预期 {expected_group_count} 组，实际仅完整映射到 {len(assigned_group_keys)} 组")
        expected_channels = ["low", "high"] if figure_type == "ihc" else [str(value).strip() for value in settings.get("channel_order", []) if str(value).strip()]
        if not expected_channels:
            expected_channels = assigned_channels
        occupied = {(group["key"], image["channel"]) for group in assigned_groups for image in group["images"]}
        for group_key in groups_to_check:
            for channel_key in expected_channels:
                if (group_key, channel_key) not in occupied:
                    add(errors, "missing_assignment_slot", f"组别 {group_key} 缺少 {channel_key} 图片")

    layout_groups: list[dict[str, Any]] = []
    try:
        if figure_type == "ihc":
            groups = build_ihc_groups(manifest, int(settings.get("group_count") or 1), settings)
            layout_groups = groups
            for group in groups:
                present = {image["channel"] for image in group["images"]}
                for channel in ("low", "high"):
                    if channel not in present:
                        add(errors, "missing_ihc_pair", f"{group['display']} 缺少 {channel} 图片")
        elif figure_type == "if":
            content = resolve_export_content(manifest, settings)
            layout_groups = content["groups"]
            if not content["groups"]:
                add(errors, "empty_figure", "没有可导出的 IF 组别")
            for group in content["groups"]:
                present = {image["channel"] for image in group["images"]}
                missing = [channel for channel in content["channel_order"] if channel not in present]
                if missing:
                    add(errors, "missing_if_channel", f"{group['display']} 缺少：{', '.join(missing)}")
    except (TypeError, ValueError) as exc:
        add(errors, "invalid_layout", str(exc))

    if isinstance(settings.get("scale_bar"), dict) and settings["scale_bar"].get("enabled"):
        cell_aspect = 1.72 if figure_type == "ihc" else 1.0
        fit_mode = str(settings.get("fit_mode") or "crop")
        seen_scale_images: set[str] = set()
        for group in layout_groups:
            for image in group.get("images", []):
                image_id = str(image.get("image_id") or "")
                if image_id in seen_scale_images:
                    continue
                seen_scale_images.add(image_id)
                try:
                    scale_bar = resolve_scale_bar(settings, image)
                    source_path = image_source_path(image)
                    _, _, _, _, visible_source_width = image_display_geometry(source_path, 0, 0, cell_aspect, 1, fit_mode)
                    if scale_bar and scale_bar["source_length_px"] > visible_source_width * 0.9:
                        add(errors, "scale_bar_too_long", f"{image['filename']} 的比例尺长度超过裁切后的可见宽度")
                except (OSError, ValueError) as exc:
                    add(errors, "invalid_scale_bar", f"{image.get('filename') or image_id}：{exc}")

    for image in ordered_manifest_images(manifest, settings):
        source_path = Path(str(image.get("source_path") or ""))
        if not source_path.is_file():
            add(errors, "missing_source", f"原图不存在或任务已过期：{image.get('filename') or image['image_id']}")
        if not image.get("sha256"):
            add(warnings, "legacy_checksum", f"{image.get('filename') or image['image_id']} 是旧任务，未记录 SHA-256")
        if not image.get("original_metadata"):
            add(warnings, "legacy_metadata", f"{image.get('filename') or image['image_id']} 是旧任务，未记录原图元数据")

    return {
        "ok": not errors,
        "errors": errors,
        "warnings": warnings,
        "error_messages": [issue["message"] for issue in errors],
        "warning_messages": [issue["message"] for issue in warnings],
        "profile": profile,
        "summary": {"errors": len(errors), "warnings": len(warnings), "images": len(manifest.get("images", []))},
    }


def build_project_package(
    manifest: dict[str, Any],
    settings: dict[str, Any],
    preflight: dict[str, Any],
    package_type: str = "project",
    final_artifact: Path | None = None,
) -> Path:
    package_type = "compliance" if package_type == "compliance" else "project"
    output_path = job_export_dir(manifest) / f"figurelab_{package_type}_{unique_export_stamp()}.zip"
    portable_images: list[dict[str, Any]] = []
    checksum_lines: list[str] = []
    used_names: set[str] = set()

    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as archive:
        for image in ordered_manifest_images(manifest, settings):
            relative = safe_path(str(image.get("relative_path") or image.get("filename") or image["image_id"]))
            archive_name = (Path("originals") / relative).as_posix()
            if archive_name in used_names:
                archive_name = (Path("originals") / relative.parent / f"{relative.stem}_{image['image_id'][:8]}{relative.suffix}").as_posix()
            used_names.add(archive_name)
            source_path = Path(image["source_path"])
            archive.write(source_path, archive_name)
            checksum = image.get("sha256") or sha256_file(source_path)
            checksum_lines.append(f"{checksum}  {archive_name}")
            portable_images.append(
                {
                    "image_id": image["image_id"],
                    "filename": image.get("filename"),
                    "relative_path": image.get("relative_path"),
                    "archive_path": archive_name,
                    "width": image.get("width"),
                    "height": image.get("height"),
                    "byte_size": image.get("byte_size") or source_path.stat().st_size,
                    "sha256": checksum,
                    "original_metadata": image.get("original_metadata") or {},
                }
            )

        final_archive_name = None
        if final_artifact:
            final_archive_name = (Path("final_figure") / final_artifact.name).as_posix()
            archive.write(final_artifact, final_archive_name)
            checksum_lines.append(f"{sha256_file(final_artifact)}  {final_archive_name}")

        project = {
            "schema": "figurelab-project/v1",
            "package_type": package_type,
            "created_at": datetime.now().isoformat(timespec="seconds"),
            "source_job_id": manifest.get("job_id"),
            "settings": {key: value for key, value in settings.items() if key != "job_id"},
            "images": portable_images,
            "final_figure": final_archive_name,
            "preflight": preflight,
        }
        archive.writestr("project.json", json.dumps(project, ensure_ascii=False, indent=2))
        archive.writestr("preflight.json", json.dumps(preflight, ensure_ascii=False, indent=2))
        archive.writestr("checksums.sha256", "\n".join(checksum_lines) + "\n")
        artifact_note = (
            "final_figure/ contains the rendered submission figure.\n"
            if final_artifact
            else "This project package may be saved before preflight passes and does not contain a rendered final figure.\n"
        )
        archive.writestr(
            "README.txt",
            f"FigureLab portable {package_type} package\n"
            "originals/ contains unmodified uploaded files.\n"
            "project.json contains layout, assignment, ROI, calibration and export settings.\n"
            f"{artifact_note}"
            "checksums.sha256 can be used to verify every packaged file.\n",
        )
    return output_path


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/wb")
@app.route("/wb/")
def wb():
    return render_template("wb.html")


@app.route("/api/upload", methods=["POST"])
def upload_folder():
    started_at = perf_counter()
    cleanup_expired_temp_data()
    files = request.files.getlist("files")
    supported_files = [file for file in files if file.filename and is_supported_image(file.filename)]
    if not supported_files:
        return jsonify({"error": "没有找到支持的 TIFF、PNG、JPG 或 BMP 图片"}), 400

    job_id = uuid.uuid4().hex
    job_dir = JOBS_DIR / job_id
    upload_dir = job_dir / "uploads"
    preview_dir = job_dir / "preview"
    upload_dir.mkdir(parents=True, exist_ok=True)

    items: list[dict[str, Any]] = []
    for index, storage in enumerate(supported_files):
        relative_path = safe_path(storage.filename)
        destination = upload_dir / relative_path
        destination.parent.mkdir(parents=True, exist_ok=True)
        if destination.exists():
            destination = destination.with_name(f"{destination.stem}_{uuid.uuid4().hex[:6]}{destination.suffix}")
        storage.save(destination)

        preview_name = f"{index:05d}_{uuid.uuid4().hex[:8]}.png"
        preview_path = preview_dir / preview_name
        try:
            image_info = prepare_image(destination, preview_path)
        except Exception as exc:
            shutil.rmtree(job_dir, ignore_errors=True)
            return jsonify({"error": f"无法读取图片 {storage.filename}：{exc}"}), 400

        items.append(
            {
                "image_id": uuid.uuid4().hex,
                "relative_path": str(relative_path).replace("\\", "/"),
                "filename": Path(storage.filename).name,
                "source_path": str(destination),
                "preview_name": preview_name,
                "width": image_info["width"],
                "height": image_info["height"],
                "sha256": sha256_file(destination),
                "byte_size": destination.stat().st_size,
                "original_metadata": image_info["metadata"],
            }
        )

    manifest = make_job_manifest(job_id, items)
    manifest_path(job_id).write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    response = public_manifest(manifest)
    response["server_processing_seconds"] = round(perf_counter() - started_at, 2)
    return jsonify(response)


@app.route("/api/jobs/<job_id>/preview/<preview_name>")
def preview(job_id: str, preview_name: str):
    if not JOB_ID_PATTERN.fullmatch(job_id) or not PREVIEW_NAME_PATTERN.fullmatch(preview_name):
        return jsonify({"error": "Preview not found"}), 404
    preview_path = JOBS_DIR / job_id / "preview" / preview_name
    if not preview_path.exists():
        return jsonify({"error": "Preview not found"}), 404
    response = send_file(preview_path)
    response.headers["Cache-Control"] = "private, no-store"
    return response


def build_export_artifact(manifest: dict[str, Any], payload: dict[str, Any]) -> Path:
    export_format = str(payload.get("export_format") or "pptx").lower()
    figure_type = str(payload.get("figure_type") or "if").lower()
    if export_format == "pptx":
        return build_ihc_pptx(manifest, payload) if figure_type == "ihc" else build_pptx(manifest, payload)
    if export_format in {"png", "jpg", "pdf", "tif", "tiff"}:
        return build_raster_export(manifest, payload, export_format)
    raise ValueError(f"Unsupported export format: {export_format}")


@app.route("/api/export", methods=["POST"])
def export_pptx():
    payload = request.get_json(force=True)
    job_id = payload.get("job_id")
    if not job_id:
        return jsonify({"error": "缺少任务编号，请重新上传"}), 400
    try:
        manifest = load_manifest(job_id)
        preflight = preflight_export(manifest, payload)
        if not preflight["ok"]:
            return jsonify({"error": preflight["error_messages"][0], "preflight": preflight}), 422
        output_path = build_export_artifact(manifest, payload)
    except Exception as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify({"download_url": f"/api/jobs/{job_id}/download/{output_path.name}", "filename": output_path.name})


@app.route("/api/preflight", methods=["POST"])
def preflight():
    payload = request.get_json(force=True)
    job_id = payload.get("job_id")
    if not job_id:
        return jsonify({"error": "缺少任务编号，请重新上传"}), 400
    try:
        manifest = load_manifest(job_id)
        report = preflight_export(manifest, payload)
    except Exception as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(report)


@app.route("/api/project/export", methods=["POST"])
def export_project():
    payload = request.get_json(force=True)
    job_id = payload.get("job_id")
    if not job_id:
        return jsonify({"error": "缺少任务编号，请重新上传"}), 400
    try:
        manifest = load_manifest(job_id)
        report = preflight_export(manifest, payload)
        output_path = build_project_package(manifest, payload, report, "project")
    except Exception as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(
        {
            "download_url": f"/api/jobs/{job_id}/download/{output_path.name}",
            "filename": output_path.name,
            "preflight": report,
        }
    )


@app.route("/api/compliance/export", methods=["POST"])
def export_compliance():
    payload = request.get_json(force=True)
    job_id = payload.get("job_id")
    if not job_id:
        return jsonify({"error": "缺少任务编号，请重新上传"}), 400
    try:
        manifest = load_manifest(job_id)
        report = preflight_export(manifest, payload)
        if not report["ok"]:
            return jsonify({"error": report["error_messages"][0], "preflight": report}), 422
        final_artifact = build_export_artifact(manifest, payload)
        output_path = build_project_package(manifest, payload, report, "compliance", final_artifact)
    except Exception as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(
        {
            "download_url": f"/api/jobs/{job_id}/download/{output_path.name}",
            "filename": output_path.name,
            "preflight": report,
        }
    )


@app.route("/api/jobs/<job_id>/download/<filename>")
def download(job_id: str, filename: str):
    if not JOB_ID_PATTERN.fullmatch(job_id) or Path(filename).name != filename or filename in {".", ".."}:
        return jsonify({"error": "Export not found"}), 404
    safe_name = Path(filename).name
    output_path = EXPORT_DIR / job_id / safe_name
    if not output_path.is_file():
        return jsonify({"error": "Export not found"}), 404
    response = send_file(output_path, as_attachment=True, download_name=safe_name)
    response.headers["Cache-Control"] = "private, no-store"
    return response


if __name__ == "__main__":
    JOBS_DIR.mkdir(parents=True, exist_ok=True)
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    port = int(os.environ.get("PORT", "5055"))
    host = "0.0.0.0" if os.environ.get("PORT") else "127.0.0.1"
    app.run(host=host, port=port, debug=False)
