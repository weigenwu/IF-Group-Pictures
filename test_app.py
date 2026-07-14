import json
import hashlib
import io
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor

import app as app_module
from app import (
    add_image_contained,
    add_scale_bar_pptx,
    add_textbox,
    build_ihc_groups,
    build_ihc_raster_pages,
    build_project_package,
    build_raster_export,
    draw_scale_bar_raster,
    image_to_rgb,
    normalized_ihc_roi,
    preflight_export,
    resolve_export_content,
    resolve_export_profile,
)


class IhcRoiTests(unittest.TestCase):
    def test_zero_and_out_of_bounds_values_are_preserved_then_clamped(self):
        self.assertEqual(normalized_ihc_roi({"ihc_roi_x": 0, "ihc_roi_y": 0}), (0.0, 0.0, 0.3, 0.28))
        self.assertEqual(normalized_ihc_roi({"ihc_roi_x": 0.9, "ihc_roi_y": 0.9, "ihc_roi_w": 0.4, "ihc_roi_h": 0.4}), (0.6, 0.6, 0.4, 0.4))

    def test_group_specific_roi_overrides_the_global_roi(self):
        settings = {
            "ihc_roi_x": 0.1,
            "ihc_rois": [{"group_key": "g2", "x": 0.2, "y": 0.3, "w": 0.4, "h": 0.5}],
        }
        self.assertEqual(normalized_ihc_roi(settings, "g1"), (0.1, 0.36, 0.3, 0.28))
        self.assertEqual(normalized_ihc_roi(settings, "g2"), (0.2, 0.3, 0.4, 0.5))


class IhcRasterSettingsTests(unittest.TestCase):
    def test_dark_contain_export_uses_the_selected_background(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            image_path = Path(temp_dir) / "square.png"
            Image.new("RGB", (60, 60), (220, 30, 30)).save(image_path)
            manifest = {
                "images": [
                    {"source_path": str(image_path)},
                    {"source_path": str(image_path)},
                ]
            }
            [page] = build_ihc_raster_pages(
                manifest,
                {
                    "group_count": 1,
                    "rows_per_slide": 1,
                    "background": "black",
                    "fit_mode": "contain",
                    "show_sample_name": False,
                    "ihc_draw_connectors": False,
                },
            )

            self.assertEqual(page.getpixel((0, 0)), (0, 0, 0))
            self.assertEqual(page.getpixel((150, page.height // 2)), (0, 0, 0))
            sampled_pixels = (
                page.getpixel((x, y))
                for y in range(0, page.height, 20)
                for x in range(0, page.width, 20)
            )
            self.assertIn((220, 30, 30), sampled_pixels)


class ExportIsolationTests(unittest.TestCase):
    def test_manual_export_filters_unchecked_channels(self):
        manifest = {"images": [{"filename": str(index)} for index in range(4)]}
        content = resolve_export_content(
            manifest,
            {
                "layout_mode": "manual",
                "group_count": 1,
                "images_per_group": 4,
                "channel_order": ["slot01", "slot03"],
            },
        )

        self.assertEqual(content["channel_order"], ["slot01", "slot03"])
        self.assertEqual([image["channel"] for image in content["groups"][0]["images"]], ["slot01", "slot03"])

    def test_exports_are_unique_and_scoped_to_the_job(self):
        job_id = "a" * 32
        with tempfile.TemporaryDirectory() as temp_dir:
            with (
                patch.object(app_module, "EXPORT_DIR", Path(temp_dir)),
                patch.object(app_module, "build_ihc_raster_pages", return_value=[Image.new("RGB", (8, 8))]),
            ):
                first = build_raster_export({"job_id": job_id}, {"figure_type": "ihc"}, "png")
                second = build_raster_export({"job_id": job_id}, {"figure_type": "ihc"}, "png")

        self.assertEqual(first.parent.name, job_id)
        self.assertNotEqual(first.name, second.name)

    def test_contained_pptx_image_gets_an_explicit_cell_background(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            image_path = Path(temp_dir) / "wide.png"
            Image.new("RGB", (80, 40), (220, 30, 30)).save(image_path)
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            add_image_contained(slide, image_path, 1, 1, 2, 2, RGBColor(0, 0, 0))

        self.assertEqual(len(slide.shapes), 2)
        self.assertEqual(tuple(slide.shapes[0].fill.fore_color.rgb), (0, 0, 0))

    def test_pptx_text_uses_arial(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        add_textbox(slide, "Marker", 1, 1, 2, 0.5, 12, RGBColor(0, 0, 0))
        self.assertEqual(slide.shapes[0].text_frame.paragraphs[0].runs[0].font.name, "Arial")


class AssignmentTests(unittest.TestCase):
    def setUp(self):
        self.manifest = {
            "images": [
                {"image_id": "a", "filename": "a.png", "relative_path": "a.png"},
                {"image_id": "b", "filename": "b.png", "relative_path": "b.png"},
                {"image_id": "c", "filename": "c.png", "relative_path": "c.png"},
                {"image_id": "d", "filename": "d.png", "relative_path": "d.png"},
            ],
            "groups": [],
            "channels": [],
        }
        self.settings = {
            "image_order": ["c", "d", "a", "b"],
            "group_order": ["g1", "g2"],
            "assignments": [
                {"image_id": "a", "group_key": "g1", "channel_key": "low"},
                {"image_id": "b", "group_key": "g1", "channel_key": "high"},
                {"image_id": "c", "group_key": "g2", "channel_key": "low"},
                {"image_id": "d", "group_key": "g2", "channel_key": "high"},
            ],
            "channel_order": ["low", "high"],
        }

    def test_assignments_drive_if_and_ihc_from_the_same_parser(self):
        content = resolve_export_content(self.manifest, self.settings)
        ihc_groups = build_ihc_groups(self.manifest, 2, self.settings)
        self.assertEqual([group["key"] for group in content["groups"]], ["g1", "g2"])
        self.assertEqual(
            [[image["image_id"] for image in group["images"]] for group in content["groups"]],
            [["a", "b"], ["c", "d"]],
        )
        self.assertEqual(
            [[image["image_id"] for image in group["images"]] for group in ihc_groups],
            [["a", "b"], ["c", "d"]],
        )


class ScientificExportTests(unittest.TestCase):
    def test_16bit_grayscale_conversion_preserves_intensity_order(self):
        source = Image.new("I;16", (5, 1))
        source.putdata([0, 1000, 10000, 30000, 65535])
        converted = image_to_rgb(source).convert("L")
        self.assertEqual(list(converted.getdata()), [0, 3, 38, 116, 255])

    def test_big_endian_16bit_grayscale_conversion_preserves_intensity_order(self):
        values = [0, 1000, 10000, 30000, 65535]
        source = Image.frombytes("I;16B", (5, 1), b"".join(value.to_bytes(2, "big") for value in values))
        converted = image_to_rgb(source).convert("L")
        self.assertEqual(list(converted.getdata()), [0, 3, 38, 116, 255])

    def test_nature_profiles_have_physical_pixel_dimensions(self):
        single = resolve_export_profile({"export_profile": "nature-single", "export_dpi": 300})
        double = resolve_export_profile({"export_profile": "nature-double", "export_dpi": 600})
        self.assertEqual(single["pixel_width"], round(8.9 / 2.54 * 300))
        self.assertEqual(double["pixel_width"], round(18.3 / 2.54 * 600))

    def test_tiff_export_is_a_real_tiff_with_dpi(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            job_id = "b" * 32
            with (
                patch.object(app_module, "EXPORT_DIR", Path(temp_dir)),
                patch.object(app_module, "build_raster_pages", return_value=[Image.new("RGB", (20, 10), (1, 2, 3))]),
            ):
                output = build_raster_export(
                    {"job_id": job_id},
                    {"figure_type": "if", "export_profile": "nature-double", "export_dpi": 600},
                    "tiff",
                )
                with Image.open(output) as image:
                    self.assertEqual(image.format, "TIFF")
                    self.assertAlmostEqual(image.info["dpi"][0], 600, delta=0.1)

    def test_calibrated_scale_bar_is_drawn_on_raster(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            source_path = Path(temp_dir) / "source.png"
            Image.new("RGB", (100, 100), (0, 0, 0)).save(source_path)
            canvas = Image.new("RGB", (200, 200), (0, 0, 0))
            settings = {
                "scale_bar": {"enabled": True, "length_um": 20, "color": "white", "position": "bottom-right", "thickness_px": 3},
                "calibration": {"pixels_per_micron": 2},
            }
            draw_scale_bar_raster(canvas, source_path, (0, 0, 200, 200), "contain", settings, 300)
            self.assertGreater(sum(pixel == (255, 255, 255) for pixel in canvas.getdata()), 50)

    def test_calibrated_scale_bar_is_editable_in_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            source_path = Path(temp_dir) / "source.png"
            Image.new("RGB", (100, 100), (0, 0, 0)).save(source_path)
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            settings = {
                "scale_bar": {"enabled": True, "length_um": 20, "color": "white", "position": "bottom-right", "thickness_px": 3},
                "calibration": {"pixels_per_micron": 2},
            }
            add_scale_bar_pptx(slide, source_path, (1, 1, 2, 2), "contain", settings, 1)
            self.assertEqual(len(slide.shapes), 2)
            self.assertEqual(slide.shapes[1].text, "20 μm")
            self.assertEqual(slide.shapes[1].text_frame.paragraphs[0].runs[0].font.name, "Arial")


class ComplianceTests(unittest.TestCase):
    def make_manifest(self, directory: Path) -> dict:
        source_path = directory / "raw.png"
        Image.new("RGB", (100, 80), (12, 34, 56)).save(source_path)
        checksum = hashlib.sha256(source_path.read_bytes()).hexdigest()
        return {
            "job_id": "c" * 32,
            "images": [
                {
                    "image_id": "image-1",
                    "filename": "raw.png",
                    "relative_path": "sample/raw.png",
                    "source_path": str(source_path),
                    "width": 100,
                    "height": 80,
                    "sha256": checksum,
                    "byte_size": source_path.stat().st_size,
                    "original_metadata": {"format": "PNG", "mode": "RGB", "frame_count": 1},
                }
            ],
            "groups": [{"key": "sample", "display": "Sample", "images": []}],
            "channels": [],
        }

    def test_preflight_rejects_invalid_scale_calibration(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            report = preflight_export(
                self.make_manifest(Path(temp_dir)),
                {
                    "figure_type": "ihc",
                    "export_format": "png",
                    "scale_bar": {"enabled": True, "length_um": 50},
                    "calibration": {},
                },
            )
        self.assertFalse(report["ok"])
        self.assertIn("invalid_scale_bar", {issue["code"] for issue in report["errors"]})

    def test_preflight_rejects_incomplete_group_mapping(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            manifest = self.make_manifest(Path(temp_dir))
            report = preflight_export(
                manifest,
                {
                    "figure_type": "if",
                    "export_format": "png",
                    "group_count": 2,
                    "channel_order": ["marker"],
                    "assignments": [{"image_id": "image-1", "group_key": "g1", "channel_key": "marker"}],
                },
            )
        self.assertFalse(report["ok"])
        self.assertIn("incomplete_group_count", {issue["code"] for issue in report["errors"]})

    def test_preflight_rejects_legacy_layout_with_missing_slots(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            manifest = self.make_manifest(Path(temp_dir))
            ihc = preflight_export(manifest, {"figure_type": "ihc", "export_format": "png", "group_count": 1})
            manual_if = preflight_export(
                manifest,
                {
                    "figure_type": "if",
                    "export_format": "png",
                    "layout_mode": "manual",
                    "group_count": 1,
                    "images_per_group": 2,
                    "channel_order": ["slot01", "slot02"],
                },
            )
        self.assertIn("missing_ihc_pair", {issue["code"] for issue in ihc["errors"]})
        self.assertIn("missing_if_channel", {issue["code"] for issue in manual_if["errors"]})

    def test_preflight_uses_crop_visible_width_for_scale_bar(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            manifest = self.make_manifest(root)
            source_path = Path(manifest["images"][0]["source_path"])
            Image.new("RGB", (1000, 100), (12, 34, 56)).save(source_path)
            manifest["images"][0].update(width=1000, height=100, sha256=hashlib.sha256(source_path.read_bytes()).hexdigest())
            report = preflight_export(
                manifest,
                {
                    "figure_type": "if",
                    "export_format": "png",
                    "group_count": 1,
                    "group_order": ["g1"],
                    "channel_order": ["marker"],
                    "fit_mode": "crop",
                    "assignments": [{"image_id": "image-1", "group_key": "g1", "channel_key": "marker"}],
                    "scale_bar": {"enabled": True, "length_um": 100},
                    "calibrations": {"default": {"pixels_per_micron": 2}},
                },
            )
        self.assertFalse(report["ok"])
        self.assertIn("scale_bar_too_long", {issue["code"] for issue in report["errors"]})

    def test_project_package_contains_sources_settings_and_checksums(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            manifest = self.make_manifest(root)
            settings = {"figure_type": "ihc", "export_profile": "nature-double", "export_dpi": 300}
            report = preflight_export(manifest, settings)
            with patch.object(app_module, "EXPORT_DIR", root / "exports"):
                package = build_project_package(manifest, settings, report)
            with zipfile.ZipFile(package) as archive:
                names = set(archive.namelist())
                project = json.loads(archive.read("project.json"))
                checksums = archive.read("checksums.sha256").decode()
            self.assertIn("originals/sample/raw.png", names)
            self.assertEqual(project["schema"], "figurelab-project/v1")
            self.assertIn(manifest["images"][0]["sha256"], checksums)


class ApiWorkflowTests(unittest.TestCase):
    @staticmethod
    def png_bytes(color: tuple[int, int, int]) -> io.BytesIO:
        stream = io.BytesIO()
        Image.new("RGB", (80, 60), color).save(stream, "PNG")
        stream.seek(0)
        return stream

    def test_upload_preflight_tiff_and_project_endpoints(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            with (
                patch.object(app_module, "JOBS_DIR", root / "jobs"),
                patch.object(app_module, "EXPORT_DIR", root / "exports"),
            ):
                client = app_module.app.test_client()
                upload = client.post(
                    "/api/upload",
                    data={
                        "files": [
                            (self.png_bytes((10, 20, 30)), "sample_low.png"),
                            (self.png_bytes((40, 50, 60)), "sample_high.png"),
                        ]
                    },
                    content_type="multipart/form-data",
                )
                self.assertEqual(upload.status_code, 200)
                job = upload.get_json()
                self.assertEqual(len(job["images"][0]["sha256"]), 64)
                self.assertEqual(job["images"][0]["original_metadata"]["format"], "PNG")

                image_ids = [image["image_id"] for image in job["images"]]
                payload = {
                    "job_id": job["job_id"],
                    "figure_type": "ihc",
                    "group_count": 1,
                    "rows_per_slide": 1,
                    "export_profile": "nature-single",
                    "export_dpi": 300,
                    "assignments": [
                        {"image_id": image_ids[0], "group_key": "g1", "channel_key": "low"},
                        {"image_id": image_ids[1], "group_key": "g1", "channel_key": "high"},
                    ],
                    "ihc_rois": [{"group_key": "g1", "x": 0.1, "y": 0.2, "w": 0.3, "h": 0.4}],
                }
                preflight = client.post("/api/preflight", json={**payload, "export_format": "tiff"})
                self.assertEqual(preflight.status_code, 200)
                self.assertTrue(preflight.get_json()["ok"])

                blocked = client.post(
                    "/api/export",
                    json={
                        **payload,
                        "export_format": "png",
                        "scale_bar": {"enabled": True, "length_um": 20},
                        "calibration": {},
                    },
                )
                self.assertEqual(blocked.status_code, 422)
                self.assertFalse(blocked.get_json()["preflight"]["ok"])

                exported = client.post("/api/export", json={**payload, "export_format": "tiff"})
                self.assertEqual(exported.status_code, 200)
                export_path = root / "exports" / job["job_id"] / exported.get_json()["filename"]
                with Image.open(export_path) as image:
                    self.assertEqual(image.format, "TIFF")
                    self.assertEqual(image.size[0], round(8.9 / 2.54 * 300))

                package = client.post(
                    "/api/project/export",
                    json={
                        **payload,
                        "package_type": "project",
                        "scale_bar": {"enabled": True, "length_um": 20},
                        "calibrations": {"by_channel": {"low": {"pixels_per_micron": 1}}},
                    },
                )
                self.assertEqual(package.status_code, 200)
                package_path = root / "exports" / job["job_id"] / package.get_json()["filename"]
                with zipfile.ZipFile(package_path) as archive:
                    self.assertIn("project.json", archive.namelist())
                    project = json.loads(archive.read("project.json"))
                    self.assertEqual(project["package_type"], "project")
                    self.assertFalse(project["preflight"]["ok"])
                    self.assertIsNone(project["final_figure"])

                compliance = client.post("/api/compliance/export", json={**payload, "export_format": "tiff"})
                self.assertEqual(compliance.status_code, 200)
                compliance_path = root / "exports" / job["job_id"] / compliance.get_json()["filename"]
                with zipfile.ZipFile(compliance_path) as archive:
                    project = json.loads(archive.read("project.json"))
                    self.assertEqual(project["package_type"], "compliance")
                    self.assertIn(project["final_figure"], archive.namelist())


if __name__ == "__main__":
    unittest.main()
